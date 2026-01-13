#!/usr/bin/env python3
"""
Samba Slides - PowerPoint Generation Library

Generate Samba TV branded presentations programmatically using python-pptx.
Implements the "Technical Warmth" design philosophy.

Usage:
    from samba_pptx import SambaPresentation

    prs = SambaPresentation()
    prs.add_title_slide("AI-Driven Media Intelligence")
    prs.add_section_divider("The Problem")
    prs.add_content_slide(
        "Traditional measurement is broken",
        ["Fragmented data sources", "No cross-platform view", "Delayed insights"]
    )
    prs.add_metrics_slide([
        {"value": "500M+", "label": "Devices", "accent": True},
        {"value": "24/7", "label": "Real-time data"},
        {"value": "99.9%", "label": "Accuracy"}
    ])

    # HUD overlays for "Tek on Img" brand element
    slide = prs.add_image_slide("path/to/photo.jpg")
    prs.add_hud_coordinates(slide, 0.85, 0.1, "X: 1040  Y: 447")
    prs.add_hud_timestamp(slide, 0.85, 0.15)
    prs.add_tracking_dot(slide, 0.5, 0.5, label="SUBJECT")

    prs.save("presentation.pptx")

Dependencies:
    pip install python-pptx Pillow
"""

import os
import sys
import shutil
from pathlib import Path
from datetime import datetime

# YAML is optional - fallback to inline manifest if not available
try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    raise ImportError("python-pptx is required. Install with: pip install python-pptx")


# =============================================================================
# ENVIRONMENT DETECTION
# =============================================================================

class SkillEnvironment:
    """Detect execution environment and resolve paths accordingly.

    Handles three contexts:
    1. Container (Anthropic's code execution) - mounted at /skills/{name}/
    2. Local development (macOS/Linux) - relative to script
    3. Custom path provided by caller
    """

    SKILL_NAME = "samba-slides"

    @classmethod
    def is_container(cls) -> bool:
        """Check if running in Anthropic's container environment."""
        # Check common container indicators
        return any([
            os.environ.get("CLAUDE_CONTAINER"),
            os.environ.get("ANTHROPIC_CONTAINER"),
            Path(f"/skills/{cls.SKILL_NAME}").exists(),
            # Container typically has no home directory or limited one
            not Path.home().exists() or str(Path.home()) == "/",
        ])

    @classmethod
    def resolve_skill_root(cls) -> Path:
        """Find the skill's root directory across execution contexts."""
        candidates = [
            # Container mount path (highest priority)
            Path(f"/skills/{cls.SKILL_NAME}"),
            # Relative to this script (local dev)
            Path(__file__).parent.parent,
            # Current working directory fallback
            Path.cwd() / ".skills" / cls.SKILL_NAME,
        ]

        for path in candidates:
            if path.exists() and (path / "assets").exists():
                return path

        # Last resort: assume script-relative
        return Path(__file__).parent.parent

    @classmethod
    def resolve_assets_dir(cls, custom_dir: Path = None) -> Path:
        """Resolve the assets directory."""
        if custom_dir and custom_dir.exists():
            return custom_dir
        return cls.resolve_skill_root() / "assets"

    @classmethod
    def should_use_fallback_fonts(cls) -> bool:
        """Determine if fallback fonts should be used.

        Returns False by default - we always reference Season font names
        in the PPTX file. python-pptx just writes font names to XML;
        it doesn't need fonts installed to reference them.

        The fonts render correctly when the PPTX is opened on a system
        with Season fonts installed. If not installed, PowerPoint/Slides
        will substitute with a similar font.

        Fonts are bundled in assets/fonts/ for users who want to install them.
        """
        return False  # Always use Season fonts - they're just name references

    @classmethod
    def get_output_dir(cls) -> Path:
        """Get appropriate output directory for generated files."""
        if cls.is_container():
            return Path("/tmp")
        return Path.cwd()


# =============================================================================
# BRAND CONSTANTS
# =============================================================================

class SambaColors:
    """Samba brand color palette."""
    VOID_BLACK = RgbColor(0x05, 0x05, 0x05)
    STARK_WHITE = RgbColor(0xFF, 0xFF, 0xFF)
    ACID_GREEN = RgbColor(0xCC, 0xFF, 0x00)
    SIGNAL_RED = RgbColor(0xFF, 0x3B, 0x30)

    # Utility grays
    DARK_GRAY = RgbColor(0x20, 0x20, 0x20)
    LIGHT_GRAY = RgbColor(0xF0, 0xF0, 0xF0)
    MID_GRAY = RgbColor(0x80, 0x80, 0x80)


class SambaFonts:
    """Samba typography specifications.

    Note: Custom fonts must be installed on the system for python-pptx to use them.
    Font files are bundled in assets/fonts/ for manual installation.
    """
    # Season Mix - Headlines, Hero Statements
    HEADLINE_FONT = "Season Mix"
    HEADLINE_LIGHT = "Season Mix Light"
    HEADLINE_MEDIUM = "Season Mix Medium"

    # Season Sans - UI, Data, Body
    BODY_FONT = "Season Sans"
    BODY_MEDIUM = "Season Sans Medium"
    BODY_SEMIBOLD = "Season Sans SemiBold"

    # Fallback fonts (if Season not installed)
    FALLBACK_HEADLINE = "Arial"
    FALLBACK_BODY = "Arial"

    # HUD/Technical fonts (monospace for data overlays)
    HUD_FONT = "Consolas"  # Monospace for technical feel
    HUD_FALLBACK = "Courier New"

    _font_cache = {}

    @classmethod
    def check_font_available(cls, font_name):
        """Check if a font is available on the system (basic check)."""
        if font_name in cls._font_cache:
            return cls._font_cache[font_name]
        # Note: python-pptx doesn't validate fonts - they're embedded by name
        # This is a placeholder for potential future font validation
        cls._font_cache[font_name] = True
        return True

    @classmethod
    def get_headline_font(cls, use_fallback=False):
        """Get headline font with fallback support."""
        if use_fallback:
            return cls.FALLBACK_HEADLINE
        return cls.HEADLINE_FONT

    @classmethod
    def get_body_font(cls, use_fallback=False):
        """Get body font with fallback support."""
        if use_fallback:
            return cls.FALLBACK_BODY
        return cls.BODY_FONT

    @classmethod
    def get_hud_font(cls):
        """Get monospace font for HUD overlays."""
        return cls.HUD_FONT


class FontManager:
    """Manages Samba brand font validation.

    Note: In container environments, fonts cannot be installed to system directories.
    The skill bundles fonts in assets/fonts/ but python-pptx only embeds font NAMES
    in the PPTX file - actual font rendering happens when the file is opened.

    For containers, we automatically use fallback fonts (Arial/Courier New).
    """

    FONT_FILES = [
        "SeasonMix-Light.ttf",
        "SeasonMix-Regular.ttf",
        "SeasonMix-Medium.ttf",
        "SeasonSans-Regular.ttf",
        "SeasonSans-Medium.ttf",
        "SeasonSans-SemiBold.ttf"
    ]

    @classmethod
    def _get_system_fonts_dir(cls) -> Path:
        """Get platform-appropriate system fonts directory."""
        if sys.platform == "darwin":
            return Path.home() / "Library/Fonts"
        elif sys.platform == "win32":
            return Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts"
        else:
            # Linux - try user fonts first
            user_fonts = Path.home() / ".local/share/fonts"
            if user_fonts.exists():
                return user_fonts
            return Path.home() / ".fonts"

    @classmethod
    def check_fonts_in_dir(cls, fonts_dir: Path) -> dict:
        """Check which Season fonts exist in a directory.

        Args:
            fonts_dir: Directory to check for font files

        Returns:
            dict: {font_file: True/False} for each font
        """
        result = {}
        for font in cls.FONT_FILES:
            font_path = fonts_dir / font
            result[font] = font_path.exists()
        return result

    @classmethod
    def check_fonts_installed(cls) -> dict:
        """Check which Season fonts are installed in system fonts directory."""
        return cls.check_fonts_in_dir(cls._get_system_fonts_dir())

    @classmethod
    def all_fonts_installed(cls) -> bool:
        """Check if all Season fonts are installed in system directory."""
        return all(cls.check_fonts_installed().values())

    @classmethod
    def fonts_bundled(cls, skill_fonts_dir: Path) -> bool:
        """Check if fonts are bundled with the skill."""
        if not skill_fonts_dir.exists():
            return False
        return all(cls.check_fonts_in_dir(skill_fonts_dir).values())

    @classmethod
    def install_fonts(cls, source_dir: Path) -> dict:
        """Copy fonts from source_dir to system fonts directory.

        Note: This is skipped in container environments where system
        directories are read-only.

        Args:
            source_dir: Path to directory containing font files

        Returns:
            dict: {font_file: True/False} for each font installation
        """
        # Skip installation in container environments
        if SkillEnvironment.is_container():
            return {font: False for font in cls.FONT_FILES}

        dest_dir = cls._get_system_fonts_dir()
        try:
            dest_dir.mkdir(parents=True, exist_ok=True)
        except (PermissionError, OSError):
            return {font: False for font in cls.FONT_FILES}

        result = {}
        for font in cls.FONT_FILES:
            source = source_dir / font
            dest = dest_dir / font

            if dest.exists():
                result[font] = True  # Already installed
            elif source.exists():
                try:
                    shutil.copy2(source, dest)
                    result[font] = True
                except (PermissionError, OSError):
                    result[font] = False
            else:
                result[font] = False  # Source not found

        return result

    @classmethod
    def ensure_fonts(cls, source_dir: Path, auto_install: bool = True) -> bool:
        """Check and optionally install fonts.

        In container environments, this always returns False (use fallback fonts).

        Args:
            source_dir: Path to directory containing font files
            auto_install: If True, install missing fonts (ignored in containers)

        Returns:
            bool: True if all fonts are available in system directory
        """
        # In containers, we can't install fonts - always use fallback
        if SkillEnvironment.is_container():
            return False

        if cls.all_fonts_installed():
            return True

        if auto_install:
            cls.install_fonts(source_dir)
            return cls.all_fonts_installed()

        return False


class LogoManager:
    """Manages context-aware logo selection for Samba presentations."""

    DEFAULT_MANIFEST = {
        "logos": {
            "default": {"dark": "samba-logo-white.png", "light": "samba-logo-black.png"},
            "internal": {"dark": "samba-logo-white.png", "light": "samba-logo-black.png"},
            "external": {"dark": "samba-logo-white.png", "light": "samba-logo-black.png"},
            "investor": {"dark": "samba-logo-white.png", "light": "samba-logo-black.png"},
            "partner": {"dark": "samba-logo-white.png", "light": "samba-logo-black.png"},
        },
        "auto_branding": {
            "internal": {
                "title": {"position": "bottom-right", "size": "small"},
                "end": {"position": "center", "size": "large"}
            },
            "external": {
                "title": {"position": "bottom-right", "size": "medium"},
                "end": {"position": "center", "size": "large"}
            },
            "investor": {
                "title": {"position": "bottom-right", "size": "medium"},
                "end": {"position": "center", "size": "large"}
            },
            "partner": {
                "end": {"position": "center", "size": "large"}
            }
        }
    }

    def __init__(self, logos_dir: Path):
        """Initialize LogoManager.

        Args:
            logos_dir: Path to logos directory
        """
        self.logos_dir = logos_dir
        self.manifest = self._load_manifest()

    def _load_manifest(self) -> dict:
        """Load logo manifest from YAML or use defaults.

        Falls back to DEFAULT_MANIFEST if:
        - YAML module not available
        - Manifest file doesn't exist
        - Manifest file can't be parsed
        """
        # If YAML isn't available, use defaults
        if not HAS_YAML:
            return self.DEFAULT_MANIFEST

        manifest_path = self.logos_dir / "logo-manifest.yaml"

        if manifest_path.exists():
            try:
                with open(manifest_path, 'r') as f:
                    return yaml.safe_load(f)
            except Exception:
                pass

        return self.DEFAULT_MANIFEST

    def get_logo_path(self, context: str, theme: str) -> Path:
        """Get appropriate logo file for context and theme.

        Args:
            context: Presentation context ('internal', 'external', 'investor', 'partner')
            theme: Theme ('dark' or 'light')

        Returns:
            Path to logo file
        """
        logos = self.manifest.get("logos", {})

        # Try context-specific logo, fallback to default
        context_logos = logos.get(context, logos.get("default", {}))
        logo_file = context_logos.get(theme, "samba-logo-white.png")

        return self.logos_dir / logo_file

    def get_auto_branding_config(self, context: str) -> dict:
        """Get auto-branding configuration for context.

        Args:
            context: Presentation context

        Returns:
            dict with slide type -> {position, size} mappings
        """
        auto_branding = self.manifest.get("auto_branding", {})
        return auto_branding.get(context, auto_branding.get("internal", {}))

    def get_available_contexts(self) -> list:
        """Get list of available presentation contexts."""
        return list(self.manifest.get("logos", {}).keys())


class SambaSizes:
    """Standard sizing for Samba presentations."""
    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins
    MARGIN = Inches(0.75)
    CONTENT_WIDTH = Inches(11.833)  # SLIDE_WIDTH - 2*MARGIN

    # Font sizes
    TITLE_HUGE = Pt(120)    # Section dividers
    TITLE_LARGE = Pt(72)    # Title slides
    HEADLINE = Pt(44)       # Content headlines
    SUBHEAD = Pt(24)        # Subtitles
    BODY = Pt(18)           # Body text
    CAPTION = Pt(14)        # Captions
    MICRO = Pt(12)          # Tiny labels, timestamps
    NANO = Pt(10)           # Technical labels

    # HUD overlay sizes
    HUD_LABEL = Pt(9)       # HUD coordinate/timestamp labels
    HUD_DOT_SIZE = Inches(0.15)  # Tracking dot diameter
    HUD_LINE_WIDTH = Pt(1)  # Wireframe/bounding box line width

    # Table sizes
    TABLE_HEADER = Pt(14)
    TABLE_CELL = Pt(12)


# =============================================================================
# HELPER CLASSES
# =============================================================================

class _ShapeGroup:
    """Logical grouping of shapes for coordinated movement.

    This is a Python-side tracking mechanism - NOT a native PPTX GroupShape.
    Shapes in this group can be moved together programmatically, but will
    appear as individual shapes in PowerPoint's UI.

    Usage:
        group = prs.create_group(slide, "my_metric")
        headline = prs.add_headline(slide, "Title", x=0.1, y=0.1)
        prs.add_to_group(group, headline)
        value = prs.add_body_text(slide, "500M+", x=0.1, y=0.2)
        prs.add_to_group(group, value)
        prs.move_group(group, dx=0.2, dy=0)  # Shift entire group right
    """

    def __init__(self, presentation, slide, name=None):
        """Initialize a shape group.

        Args:
            presentation: SambaPresentation instance (for _to_inches access)
            slide: Slide object
            name: Optional name for the group
        """
        self.presentation = presentation
        self.slide = slide
        self.name = name
        self.shapes = []

    def add_shape(self, shape):
        """Register a shape with this group.

        Args:
            shape: Shape object to add
        """
        if shape and shape not in self.shapes:
            self.shapes.append(shape)

    def remove_shape(self, shape):
        """Remove a shape from this group.

        Args:
            shape: Shape object to remove
        """
        if shape in self.shapes:
            self.shapes.remove(shape)

    def move(self, dx, dy, unit="ratio"):
        """Move all shapes by delta x/y.

        Args:
            dx: Horizontal delta (positive = right)
            dy: Vertical delta (positive = down)
            unit: "ratio", "inches", or "px"

        Note:
            Shape positions in python-pptx are stored in EMUs (English Metric Units).
            Inches objects auto-convert when added to EMU values.
        """
        dx_emu = Emu(self.presentation._to_inches(dx, "width", unit))
        dy_emu = Emu(self.presentation._to_inches(dy, "height", unit))

        for shape in self.shapes:
            try:
                shape.left = shape.left + dx_emu
                shape.top = shape.top + dy_emu
            except AttributeError:
                # Some shapes (like connectors) may not have left/top
                pass

    def get_bounds(self):
        """Get the bounding box of all shapes in the group.

        Returns:
            Tuple of (left, top, right, bottom) in EMUs, or None if empty
        """
        if not self.shapes:
            return None

        min_left = min(s.left for s in self.shapes if hasattr(s, 'left'))
        min_top = min(s.top for s in self.shapes if hasattr(s, 'top'))
        max_right = max(s.left + s.width for s in self.shapes
                       if hasattr(s, 'left') and hasattr(s, 'width'))
        max_bottom = max(s.top + s.height for s in self.shapes
                        if hasattr(s, 'top') and hasattr(s, 'height'))

        return min_left, min_top, max_right, max_bottom


# =============================================================================
# MAIN CLASS
# =============================================================================

class SambaPresentation:
    """Generate Samba TV branded PowerPoint presentations.

    This class provides two API layers:

    **Layer 2: Brand Primitives** (NEW - Flexible positioning)
        - add_slide() - Create a blank themed slide
        - add_headline() - Add positioned headline text
        - add_body_text() - Add positioned body text
        - add_bullets() - Add bullet list
        - add_metric() - Add value/label metric pair
        - add_rectangle() - Add colored rectangle
        - add_accent_bar() - Add brand accent bar
        - add_circle() - Add circle shape
        - add_line() - Add line shape
        - add_image() - Add positioned image
        - add_video() - Add embedded video
        - set_background_image() - Set slide background
        - add_image_grid() - Add grid of images
        - add_logo() - Add logo with flexible positioning

    **Layer 3: Convenience Templates** (Existing - Use primitives internally)
        - add_title_slide() - Pre-designed title slide
        - add_content_slide() - Pre-designed content slide
        - add_metrics_slide() - Pre-designed metrics slide
        - etc.

    Coordinates can be specified as:
        - Ratios (default): 0.0-1.0 representing percentage of slide dimensions
        - Inches: Use unit="inches" parameter
        - Pixels: Use unit="px" parameter
    """

    # Valid presentation contexts
    VALID_CONTEXTS = ["internal", "external", "investor", "partner"]

    def __init__(
        self,
        theme="dark",
        context="internal",
        auto_branding=True,
        auto_install_fonts=True,
        assets_dir=None,
        use_fallback_fonts=None  # Changed: None means "auto-detect"
    ):
        """
        Initialize a new Samba presentation.

        Args:
            theme: "dark" (Void Black bg) or "light" (Stark White bg)
            context: Presentation context - determines logo selection and branding.
                    Options: "internal" (team meetings), "external" (client-facing),
                    "investor" (board/investor), "partner" (co-branded materials)
            auto_branding: If True, automatically add logo to title and end slides
            auto_install_fonts: If True, install Season fonts if not present
                               (ignored in container environments)
            assets_dir: Custom path to assets folder (logos, fonts). If None,
                       uses SkillEnvironment to resolve automatically.
            use_fallback_fonts: If True, use Arial instead of Season fonts.
                               If None (default), auto-detect based on environment
                               (containers always use fallback).
        """
        # Resolve assets directory using SkillEnvironment
        if assets_dir:
            self.assets_dir = Path(assets_dir)
        else:
            self.assets_dir = SkillEnvironment.resolve_assets_dir()

        self.logos_dir = self.assets_dir / "logos"
        self.fonts_dir = self.assets_dir / "fonts"

        # Auto-detect fallback fonts setting
        # In containers, always use fallback (can't install system fonts)
        if use_fallback_fonts is None:
            use_fallback_fonts = SkillEnvironment.should_use_fallback_fonts()

        # Font pre-flight: auto-install Season fonts if needed (skipped in containers)
        if auto_install_fonts and not use_fallback_fonts:
            FontManager.ensure_fonts(self.fonts_dir, auto_install=True)

        # Initialize presentation
        self.prs = Presentation()
        self.prs.slide_width = SambaSizes.SLIDE_WIDTH
        self.prs.slide_height = SambaSizes.SLIDE_HEIGHT
        self.theme = theme
        self.use_fallback_fonts = use_fallback_fonts

        # Context and branding settings
        self.context = context if context in self.VALID_CONTEXTS else "internal"
        self.auto_branding = auto_branding

        # Initialize logo manager
        self.logo_manager = LogoManager(self.logos_dir)

    def _get_headline_font(self):
        """Get the appropriate headline font."""
        return SambaFonts.get_headline_font(self.use_fallback_fonts)

    def _get_body_font(self):
        """Get the appropriate body font."""
        return SambaFonts.get_body_font(self.use_fallback_fonts)

    def _get_colors(self, theme=None):
        """Get background and text colors for theme."""
        theme = theme or self.theme
        if theme == "dark":
            return SambaColors.VOID_BLACK, SambaColors.STARK_WHITE
        else:
            return SambaColors.STARK_WHITE, SambaColors.VOID_BLACK

    # =========================================================================
    # LAYER 2: BRAND PRIMITIVES - HELPER METHODS
    # =========================================================================

    def _to_inches(self, value, dimension="width", unit="ratio"):
        """Convert value to Inches based on unit type.

        Args:
            value: The numeric value to convert
            dimension: "width" or "height" (used for ratio conversion)
            unit: "ratio" (0.0-1.0), "inches", or "px"

        Returns:
            Inches object
        """
        if unit == "inches":
            return Inches(value)
        elif unit == "px":
            # Assuming 96 DPI (standard screen), 1 inch = 96 px
            return Inches(value / 96)
        else:  # ratio (default)
            if dimension == "width":
                return Inches(value * 13.333)
            else:  # height
                return Inches(value * 7.5)

    def _resolve_color(self, color, theme=None):
        """Resolve color parameter to RGBColor.

        Args:
            color: One of:
                - "text" - auto-selects based on theme (white on dark, black on light)
                - "accent" - Acid Green
                - "white" - Stark White
                - "black" - Void Black
                - "gray" - Mid Gray
                - "#RRGGBB" - Hex string
                - RGBColor object - used directly
            theme: Optional theme override for "text" resolution

        Returns:
            RGBColor object
        """
        if isinstance(color, RgbColor):
            return color  # Already an RGBColor, use directly
        elif color == "text":
            _, text_color = self._get_colors(theme)
            return text_color
        elif color == "accent":
            return SambaColors.ACID_GREEN
        elif color == "white":
            return SambaColors.STARK_WHITE
        elif color == "black":
            return SambaColors.VOID_BLACK
        elif color == "gray":
            return SambaColors.MID_GRAY
        elif isinstance(color, str) and color.startswith("#"):
            return RgbColor.from_string(color[1:])
        return SambaColors.STARK_WHITE  # Fallback default

    def _resolve_size(self, size_name):
        """Convert size name to Pt value.

        Args:
            size_name: "huge", "large", "medium", "subhead", "body", "caption", "micro"

        Returns:
            Pt object with font size
        """
        sizes = {
            "huge": SambaSizes.TITLE_HUGE,      # 120pt
            "large": SambaSizes.TITLE_LARGE,    # 72pt
            "medium": SambaSizes.HEADLINE,      # 44pt
            "subhead": SambaSizes.SUBHEAD,      # 24pt
            "body": SambaSizes.BODY,            # 18pt
            "caption": SambaSizes.CAPTION,      # 14pt
            "micro": SambaSizes.MICRO,          # 12pt
        }
        return sizes.get(size_name, SambaSizes.BODY)

    def _resolve_alignment(self, align):
        """Convert alignment string to PP_ALIGN enum.

        Args:
            align: "left", "center", or "right"

        Returns:
            PP_ALIGN enum value
        """
        alignments = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }
        return alignments.get(align, PP_ALIGN.LEFT)

    # =========================================================================
    # LAYER 2: BRAND PRIMITIVES - SLIDE & TEXT
    # =========================================================================

    def add_slide(self, theme=None):
        """Create a new blank slide with themed background.

        This is the primary method for creating slides when using primitives.
        Use this instead of _add_blank_slide() for the public API.

        Args:
            theme: "dark" or "light" (default: uses instance theme)

        Returns:
            Slide object ready for adding primitives
        """
        return self._add_blank_slide(theme)

    def add_headline(self, slide, text, x=0.05, y=0.1, width=0.9, height=None,
                     size="large", color="text", align="left", unit="ratio", theme=None):
        """Add a branded headline to a slide.

        Args:
            slide: Slide object to add headline to
            text: Headline text
            x: X position (default 0.05 ratio)
            y: Y position (default 0.1 ratio)
            width: Width (default 0.9 ratio)
            height: Height (default auto-calculated)
            size: "huge", "large", "medium" - maps to SambaSizes
            color: "text", "accent", "white", "black", hex string, or RGBColor
            align: "left", "center", "right"
            unit: "ratio", "inches", or "px"
            theme: Override theme for color resolution

        Returns:
            TextBox shape object
        """
        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)
        h = self._to_inches(height, "height", unit) if height else Inches(1.5)

        text_box = slide.shapes.add_textbox(left, top, w, h)
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = self._resolve_size(size)
        p.font.color.rgb = self._resolve_color(color, theme or self.theme)
        p.font.bold = False
        p.font.name = self._get_headline_font()
        p.alignment = self._resolve_alignment(align)

        return text_box

    def add_body_text(self, slide, text, x=0.05, y=0.25, width=0.9, height=0.5,
                      size="body", color="text", align="left", wrap=True,
                      unit="ratio", theme=None):
        """Add body text to a slide.

        Args:
            slide: Slide object to add text to
            text: Body text content
            x: X position (default 0.05 ratio)
            y: Y position (default 0.25 ratio)
            width: Width (default 0.9 ratio)
            height: Height (default 0.5 ratio)
            size: "body", "caption", "micro", "subhead" - maps to SambaSizes
            color: "text", "accent", "white", "black", hex string, or RGBColor
            align: "left", "center", "right"
            wrap: Enable word wrap (default True)
            unit: "ratio", "inches", or "px"
            theme: Override theme for color resolution

        Returns:
            TextBox shape object
        """
        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)
        h = self._to_inches(height, "height", unit)

        text_box = slide.shapes.add_textbox(left, top, w, h)
        tf = text_box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = self._resolve_size(size)
        p.font.color.rgb = self._resolve_color(color, theme or self.theme)
        p.font.name = self._get_body_font()
        p.alignment = self._resolve_alignment(align)

        return text_box

    def add_bullets(self, slide, items, x=0.05, y=0.25, width=0.9, height=0.6,
                    size="body", color="text", bullet_color=None, spacing=12,
                    unit="ratio", theme=None):
        """Add a bullet list to a slide.

        Args:
            slide: Slide object to add bullets to
            items: List of bullet point strings
            x: X position (default 0.05 ratio)
            y: Y position (default 0.25 ratio)
            width: Width (default 0.9 ratio)
            height: Height (default 0.6 ratio)
            size: Font size name (default "body")
            color: Text color (default "text")
            bullet_color: Bullet marker color (default: same as text)
            spacing: Space after each item in points (default 12)
            unit: "ratio", "inches", or "px"
            theme: Override theme for color resolution

        Returns:
            TextBox shape object, or None if items is empty
        """
        # Handle empty list gracefully
        if not items:
            return None

        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)
        h = self._to_inches(height, "height", unit)

        text_color = self._resolve_color(color, theme or self.theme)

        text_box = slide.shapes.add_textbox(left, top, w, h)
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = self._resolve_size(size)
            p.font.color.rgb = text_color
            p.font.name = self._get_body_font()
            p.space_after = Pt(spacing)

        return text_box

    def add_rectangle(self, slide, x=0.1, y=0.1, width=0.3, height=0.2,
                      fill="accent", opacity=1.0, border=False, border_color="white",
                      unit="ratio"):
        """Add a colored rectangle to a slide.

        Args:
            slide: Slide object to add rectangle to
            x: X position (default 0.1 ratio)
            y: Y position (default 0.1 ratio)
            width: Width (default 0.3 ratio)
            height: Height (default 0.2 ratio)
            fill: Fill color ("accent", "black", "white", hex, or RGBColor)
            opacity: Fill opacity 0.0-1.0 (Note: limited PPTX support)
            border: Show border (default False)
            border_color: Border color if border=True
            unit: "ratio", "inches", or "px"

        Returns:
            Shape object
        """
        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)
        h = self._to_inches(height, "height", unit)

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = self._resolve_color(fill)

        if border:
            rect.line.color.rgb = self._resolve_color(border_color)
            rect.line.width = Pt(1)
        else:
            rect.line.fill.background()

        return rect

    def add_metric(self, slide, value, label, x=0.1, y=0.3, width=0.25,
                   accent=False, value_size="huge", label_size="micro",
                   color="text", spacing=None, unit="ratio", theme=None):
        """Add a metric display (large value + small label).

        Args:
            slide: Slide object to add metric to
            value: The metric value (e.g., "500M+", "99.9%")
            label: The label below the value (e.g., "DEVICES")
            x: X position (default 0.1 ratio)
            y: Y position (default 0.3 ratio)
            width: Width for the metric block (default 0.25 ratio)
            accent: If True, value is displayed in Acid Green
            value_size: Font size for value (default "huge" = 120pt)
            label_size: Font size for label (default "micro" = 12pt)
            color: Base text color when not accented (default "text")
            spacing: Space between value and label in inches (default: auto-calculated from value_size)
            unit: "ratio", "inches", or "px"
            theme: Override theme for color resolution

        Returns:
            Tuple of (value_textbox, label_textbox)
        """
        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)

        text_color = self._resolve_color(color, theme or self.theme)
        value_color = SambaColors.ACID_GREEN if accent else text_color

        # Calculate value box height based on font size
        value_font_size = self._resolve_size(value_size)
        # Approximate: 1pt ~= 1/72 inch, add 50% padding for line height
        value_height_inches = (value_font_size.pt / 72) * 1.5
        value_height = Inches(max(value_height_inches, 0.5))  # Minimum 0.5 inches

        # Large value
        val_box = slide.shapes.add_textbox(left, top, w, value_height)
        vf = val_box.text_frame
        vp = vf.paragraphs[0]
        vp.text = str(value)
        vp.font.size = value_font_size
        vp.font.color.rgb = value_color
        vp.font.bold = False
        vp.font.name = self._get_headline_font()

        # Calculate label position: below value with dynamic spacing
        if spacing is not None:
            label_offset = Inches(spacing)
        else:
            # Auto-spacing: roughly 1.1x the value height
            label_offset = Inches(value_height_inches * 1.1)

        # Small label below
        lab_box = slide.shapes.add_textbox(left, top + label_offset, w, Inches(0.5))
        lf = lab_box.text_frame
        lp = lf.paragraphs[0]
        lp.text = label.upper()
        lp.font.size = self._resolve_size(label_size)
        lp.font.color.rgb = text_color
        lp.font.name = self._get_body_font()

        return val_box, lab_box

    def add_accent_bar(self, slide, position="left", thickness=0.015, color="accent",
                       length=1.0, unit="ratio"):
        """Add a brand accent bar to slide edge.

        Args:
            slide: Slide object to add bar to
            position: "left", "right", "top", "bottom"
            thickness: Bar thickness (default 0.015 ratio ~= 2px)
            color: Bar color (default "accent" = Acid Green)
            length: Bar length as ratio of edge (default 1.0 = full edge)
            unit: "ratio", "inches", or "px"

        Returns:
            Shape object
        """
        bar_color = self._resolve_color(color)

        if position == "left":
            x, y = 0, (1 - length) / 2
            w, h = thickness, length
        elif position == "right":
            x, y = 1 - thickness, (1 - length) / 2
            w, h = thickness, length
        elif position == "top":
            x, y = (1 - length) / 2, 0
            w, h = length, thickness
        elif position == "bottom":
            x, y = (1 - length) / 2, 1 - thickness
            w, h = length, thickness
        else:
            x, y, w, h = 0, 0, thickness, 1  # Default to left

        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        width = self._to_inches(w, "width", unit)
        height = self._to_inches(h, "height", unit)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        return bar

    def add_circle(self, slide, x=0.5, y=0.5, size=0.1, fill="accent",
                   border=False, border_color="white", unit="ratio"):
        """Add a circle shape to a slide.

        Args:
            slide: Slide object to add circle to
            x: X position of center (default 0.5 ratio)
            y: Y position of center (default 0.5 ratio)
            size: Diameter (default 0.1 ratio)
            fill: Fill color (default "accent")
            border: Show border (default False)
            border_color: Border color if border=True
            unit: "ratio", "inches", or "px"

        Returns:
            Shape object

        Note:
            For ratio coordinates, size is calculated using width dimension
            to maintain a perfect circle (not ellipse).
        """
        # Calculate diameter using width dimension for consistent aspect ratio
        # This ensures circles remain circular even with ratio coordinates
        diameter = self._to_inches(size, "width", unit)
        half_size = diameter / 2

        # Calculate top-left from center position
        left = self._to_inches(x, "width", unit) - half_size
        top = self._to_inches(y, "height", unit) - half_size

        # Use same diameter for both width and height to create a true circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
        circle.fill.solid()
        circle.fill.fore_color.rgb = self._resolve_color(fill)

        if border:
            circle.line.color.rgb = self._resolve_color(border_color)
            circle.line.width = Pt(1)
        else:
            circle.line.fill.background()

        return circle

    def add_line(self, slide, start=(0.1, 0.5), end=(0.9, 0.5),
                 color="accent", thickness=2, unit="ratio"):
        """Add a line to a slide.

        Args:
            slide: Slide object to add line to
            start: (x, y) tuple for line start (default (0.1, 0.5))
            end: (x, y) tuple for line end (default (0.9, 0.5))
            color: Line color (default "accent")
            thickness: Line thickness in points (default 2)
            unit: "ratio", "inches", or "px"

        Returns:
            Connector shape object
        """
        from pptx.enum.shapes import MSO_CONNECTOR

        x1 = self._to_inches(start[0], "width", unit)
        y1 = self._to_inches(start[1], "height", unit)
        x2 = self._to_inches(end[0], "width", unit)
        y2 = self._to_inches(end[1], "height", unit)

        # python-pptx uses add_connector for lines
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2
        )
        connector.line.color.rgb = self._resolve_color(color)
        connector.line.width = Pt(thickness)

        return connector

    # =========================================================================
    # LAYER 2: BRAND PRIMITIVES - MEDIA
    # =========================================================================

    def add_image(self, slide, path, x=0.1, y=0.1, width=0.4, height=None,
                  fit="contain", unit="ratio"):
        """Add an image to a slide with flexible positioning.

        Args:
            slide: Slide object to add image to
            path: Path to image file
            x: X position (default 0.1 ratio)
            y: Y position (default 0.1 ratio)
            width: Width (default 0.4 ratio)
            height: Height (default: auto-calculated to maintain aspect ratio)
            fit: "contain" (maintain aspect ratio) or "fill" (stretch to fit)
            unit: "ratio", "inches", or "px"

        Returns:
            Picture shape object, or None if file not found
        """
        if not path or not os.path.exists(path):
            return None

        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)

        if height and fit == "fill":
            h = self._to_inches(height, "height", unit)
            picture = slide.shapes.add_picture(path, left, top, width=w, height=h)
        else:
            # Let python-pptx maintain aspect ratio
            picture = slide.shapes.add_picture(path, left, top, width=w)

        return picture

    def add_video(self, slide, path, x=0.1, y=0.1, width=0.8, height=0.6,
                  poster=None, unit="ratio"):
        """Add a video to a slide with poster frame.

        Args:
            slide: Slide object to add video to
            path: Path to video file (mp4, wmv, etc.)
            x: X position (default 0.1 ratio)
            y: Y position (default 0.1 ratio)
            width: Width (default 0.8 ratio)
            height: Height (default 0.6 ratio)
            poster: Path to poster/thumbnail image (shown before play)
            unit: "ratio", "inches", or "px"

        Returns:
            Movie shape object, or None if file not found

        Note:
            Video playback requires PowerPoint; Keynote has limited support.
        """
        if not path or not os.path.exists(path):
            return None

        left = self._to_inches(x, "width", unit)
        top = self._to_inches(y, "height", unit)
        w = self._to_inches(width, "width", unit)
        h = self._to_inches(height, "height", unit)

        # Determine mime type from extension
        ext = os.path.splitext(path)[1].lower()
        mime_types = {
            '.mp4': 'video/mp4',
            '.wmv': 'video/x-ms-wmv',
            '.avi': 'video/x-msvideo',
            '.mov': 'video/quicktime',
        }
        mime_type = mime_types.get(ext, 'video/mp4')

        movie = slide.shapes.add_movie(
            path, left, top, w, h,
            poster_frame_image=poster,
            mime_type=mime_type
        )

        return movie

    def set_background_image(self, slide, path, opacity=1.0, overlay=None):
        """Set a background image for a slide.

        Args:
            slide: Slide object to set background for
            path: Path to image file
            opacity: Image opacity 0.0-1.0 (Note: limited PPTX support)
            overlay: Optional overlay color to darken image (e.g., "black")

        Returns:
            Picture shape object, or None if file not found

        Note:
            This adds a full-slide image behind other content.
            For true background fill, use slide.background (not supported in python-pptx).
        """
        if not path or not os.path.exists(path):
            return None

        # Add image at full slide size, positioned at (0, 0)
        picture = slide.shapes.add_picture(
            path,
            Inches(0), Inches(0),
            width=SambaSizes.SLIDE_WIDTH,
            height=SambaSizes.SLIDE_HEIGHT
        )

        # Move to back (z-order)
        # python-pptx doesn't have direct z-order API, but shapes added first are behind
        # We manipulate the XML to move the image behind other shapes
        spTree = slide.shapes._spTree
        sp = picture._element
        spTree.remove(sp)

        # Find the first shape element (sp or pic) to insert before
        # spTree typically starts with nvGrpSpPr and grpSpPr, then shapes
        first_shape_idx = 2  # Default: after nvGrpSpPr and grpSpPr
        for i, child in enumerate(spTree):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('sp', 'pic', 'graphicFrame', 'cxnSp'):
                first_shape_idx = i
                break
        spTree.insert(first_shape_idx, sp)

        # Add semi-transparent overlay if requested
        if overlay:
            overlay_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                SambaSizes.SLIDE_WIDTH, SambaSizes.SLIDE_HEIGHT
            )
            overlay_rect.fill.solid()
            overlay_rect.fill.fore_color.rgb = self._resolve_color(overlay)
            overlay_rect.line.fill.background()

            # Move overlay behind other shapes but in front of background image
            sp_overlay = overlay_rect._element
            spTree.remove(sp_overlay)
            spTree.insert(first_shape_idx + 1, sp_overlay)

        return picture

    def add_image_grid(self, slide, images, x=0.05, y=0.15, width=0.9, height=0.7,
                       cols=3, gap=0.02, fit="contain", unit="ratio"):
        """Add a grid of images to a slide.

        Args:
            slide: Slide object to add grid to
            images: List of image paths
            x: X position of grid (default 0.05 ratio)
            y: Y position of grid (default 0.15 ratio)
            width: Total width of grid (default 0.9 ratio)
            height: Total height of grid (default 0.7 ratio)
            cols: Number of columns (default 3)
            gap: Gap between images as ratio (default 0.02)
            fit: "contain" (maintain aspect ratio, default) or "fill" (stretch to fit cell)
            unit: "ratio", "inches", or "px"

        Returns:
            List of Picture shape objects
        """
        if not images:
            return []

        num_images = len(images)
        rows = (num_images + cols - 1) // cols  # Ceiling division

        # Calculate cell dimensions
        total_gap_w = gap * (cols - 1)
        total_gap_h = gap * (rows - 1)
        cell_w = (width - total_gap_w) / cols
        cell_h = (height - total_gap_h) / rows

        pictures = []
        for i, img_path in enumerate(images):
            if not img_path or not os.path.exists(img_path):
                continue

            row = i // cols
            col = i % cols

            img_x = x + col * (cell_w + gap)
            img_y = y + row * (cell_h + gap)

            pic = self.add_image(slide, img_path, img_x, img_y, cell_w, cell_h,
                                fit=fit, unit=unit)
            if pic:
                pictures.append(pic)

        return pictures

    # =========================================================================
    # LAYER 2: BRAND PRIMITIVES - ELEMENT GROUPING (OPTIONAL)
    # =========================================================================

    def create_group(self, slide, name=None):
        """Create a logical group for coordinated shape movement.

        This is a Python-side tracking mechanism, NOT a native PPTX GroupShape.
        Shapes in this group can be moved together programmatically, but will
        appear as individual shapes in PowerPoint's UI.

        Args:
            slide: Slide object
            name: Optional name for the group

        Returns:
            _ShapeGroup tracker object
        """
        return _ShapeGroup(self, slide, name)

    def add_to_group(self, group, shape):
        """Add a shape to a logical group.

        Args:
            group: _ShapeGroup object from create_group()
            shape: Shape object to add to the group
        """
        group.add_shape(shape)

    def move_group(self, group, dx, dy, unit="ratio"):
        """Move all shapes in a group by delta.

        Args:
            group: _ShapeGroup object
            dx: Horizontal delta (positive = right)
            dy: Vertical delta (positive = down)
            unit: "ratio", "inches", or "px"
        """
        group.move(dx, dy, unit)

    def _add_blank_slide(self, theme=None):
        """Add a blank slide with themed background."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        bg_color, _ = self._get_colors(theme)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        return slide

    # =========================================================================
    # SLIDE TYPES
    # =========================================================================

    def add_title_slide(self, title, subtitle=None, theme=None, include_logo=None):
        """
        Add a title slide with massive typography.

        Args:
            title: Main title text (rendered huge)
            subtitle: Optional subtitle/section label (rendered small)
            theme: Override default theme for this slide
            include_logo: Whether to add logo (default: uses auto_branding setting)
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        # Massive headline
        title_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, Inches(0.5),
            SambaSizes.CONTENT_WIDTH, Inches(4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = SambaSizes.TITLE_LARGE
        p.font.color.rgb = text_color
        p.font.bold = False
        p.font.name = self._get_headline_font()

        # Small subtitle at bottom
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                SambaSizes.MARGIN, Inches(6.5),
                Inches(6), Inches(0.5)
            )
            sf = sub_box.text_frame
            sp = sf.paragraphs[0]
            sp.text = subtitle
            sp.font.size = SambaSizes.BODY
            sp.font.color.rgb = text_color
            sp.font.name = self._get_body_font()

        # Auto-branding: add logo based on context settings
        if include_logo is None:
            include_logo = self.auto_branding

        if include_logo:
            branding_config = self.logo_manager.get_auto_branding_config(self.context)
            title_config = branding_config.get("title")
            if title_config:
                self.add_logo_to_slide(
                    slide,
                    position=title_config.get("position", "bottom-right"),
                    size=title_config.get("size", "small"),
                    theme=theme
                )

        return slide

    def add_section_divider(self, text, theme=None, align="center"):
        """
        Add a section divider with extreme-scale typography.

        Args:
            text: Single word or short phrase
            theme: Override default theme
            align: Text alignment - "center" (default), "left", or "right"
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        text_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, Inches(2),
            SambaSizes.CONTENT_WIDTH, Inches(3.5)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = SambaSizes.TITLE_HUGE
        p.font.color.rgb = text_color
        p.font.bold = False
        p.font.name = self._get_headline_font()

        # Alignment mapping
        alignments = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT
        }
        p.alignment = alignments.get(align, PP_ALIGN.CENTER)

        return slide

    def add_content_slide(self, headline, points, theme=None, image_path=None):
        """
        Add a split-layout content slide (50/50).

        Args:
            headline: Section headline
            points: List of bullet points
            theme: Override default theme
            image_path: Optional path to image for right side
        """
        slide = self._add_blank_slide(theme)
        bg_color, text_color = self._get_colors(theme)

        # LEFT SIDE: Text content (50%)
        # Headline
        head_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, SambaSizes.MARGIN,
            Inches(5.5), Inches(1.5)
        )
        hf = head_box.text_frame
        hp = hf.paragraphs[0]
        hp.text = headline
        hp.font.size = SambaSizes.HEADLINE
        hp.font.color.rgb = text_color
        hp.font.bold = False
        hp.font.name = self._get_headline_font()

        # Body bullets
        body_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, Inches(2.5),
            Inches(5.5), Inches(4)
        )
        bf = body_box.text_frame
        bf.word_wrap = True

        for i, point in enumerate(points):
            if i == 0:
                p = bf.paragraphs[0]
            else:
                p = bf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = SambaSizes.BODY
            p.font.color.rgb = text_color
            p.font.name = self._get_body_font()
            p.space_after = Pt(12)

        # RIGHT SIDE: Image or placeholder (50%)
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(7), SambaSizes.MARGIN,
                width=Inches(5.5)
            )
        else:
            # Add placeholder rectangle
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7), SambaSizes.MARGIN,
                Inches(5.5), Inches(6)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = (
                SambaColors.DARK_GRAY if (theme or self.theme) == "dark"
                else SambaColors.LIGHT_GRAY
            )
            placeholder.line.fill.background()

        return slide

    def add_metrics_slide(self, metrics, theme=None):
        """
        Add a metrics/data slide with large numbers.

        Args:
            metrics: List of dicts with 'value', 'label', 'accent' (optional bool)
            theme: Override default theme

        Example:
            [
                {"value": "500M+", "label": "Devices", "accent": True},
                {"value": "24/7", "label": "Real-time"},
                {"value": "99.9%", "label": "Accuracy"}
            ]
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        num_metrics = len(metrics)
        width_per = 12 / num_metrics

        for i, metric in enumerate(metrics):
            x_pos = 0.75 + (i * width_per)

            # Large value
            val_box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(2.5),
                Inches(width_per - 0.5), Inches(2)
            )
            vf = val_box.text_frame
            vp = vf.paragraphs[0]
            vp.text = str(metric['value'])
            vp.font.size = Pt(96)
            vp.font.color.rgb = (
                SambaColors.ACID_GREEN if metric.get('accent')
                else text_color
            )
            vp.font.bold = False
            vp.font.name = self._get_headline_font()

            # Tiny label
            lab_box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(4.5),
                Inches(width_per - 0.5), Inches(0.5)
            )
            lf = lab_box.text_frame
            lp = lf.paragraphs[0]
            lp.text = metric['label'].upper()
            lp.font.size = SambaSizes.MICRO
            lp.font.color.rgb = text_color
            lp.font.name = self._get_body_font()

        return slide

    def add_quote_slide(self, quote, attribution=None, theme=None):
        """
        Add a quote slide with large italic text.

        Args:
            quote: The quote text
            attribution: Optional source/speaker
            theme: Override default theme
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        # Quote
        quote_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2),
            Inches(10.3), Inches(3)
        )
        qf = quote_box.text_frame
        qf.word_wrap = True
        qp = qf.paragraphs[0]
        qp.text = f'"{quote}"'
        qp.font.size = Pt(36)
        qp.font.color.rgb = text_color
        qp.font.italic = True
        qp.font.name = self._get_headline_font()

        # Attribution
        if attribution:
            attr_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(5.5),
                Inches(10), Inches(0.5)
            )
            af = attr_box.text_frame
            ap = af.paragraphs[0]
            ap.text = f"— {attribution}"
            ap.font.size = SambaSizes.BODY
            ap.font.color.rgb = SambaColors.MID_GRAY
            ap.font.name = self._get_body_font()

        return slide

    def add_end_slide(self, include_logo=None, text="Thank You", theme=None):
        """
        Add an end/thank you slide with optional logo.

        Args:
            include_logo: Whether to include Samba logo (default: uses auto_branding setting)
            text: Text to display (default "Thank You")
            theme: Override default theme
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        # Centered text
        text_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, Inches(2.5),
            SambaSizes.CONTENT_WIDTH, Inches(2)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = SambaSizes.TITLE_LARGE
        p.font.color.rgb = text_color
        p.font.name = self._get_headline_font()
        p.alignment = PP_ALIGN.CENTER

        # Determine if logo should be added
        if include_logo is None:
            include_logo = self.auto_branding

        # Add logo using LogoManager
        if include_logo:
            branding_config = self.logo_manager.get_auto_branding_config(self.context)
            end_config = branding_config.get("end", {"position": "center", "size": "large"})

            self.add_logo_to_slide(
                slide,
                position=end_config.get("position", "center"),
                size=end_config.get("size", "large"),
                theme=theme
            )

        return slide

    # =========================================================================
    # NEW SLIDE TYPES
    # =========================================================================

    def add_image_slide(self, image_path, caption=None, theme=None):
        """
        Add a full-bleed image slide with optional caption.

        Args:
            image_path: Path to the image file
            caption: Optional caption text (appears at bottom)
            theme: Override default theme
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        if image_path and os.path.exists(image_path):
            # Full-bleed image (with small margin)
            slide.shapes.add_picture(
                image_path,
                Inches(0.25), Inches(0.25),
                width=Inches(12.833),
                height=Inches(6.5) if not caption else Inches(6)
            )

        if caption:
            cap_box = slide.shapes.add_textbox(
                SambaSizes.MARGIN, Inches(6.75),
                SambaSizes.CONTENT_WIDTH, Inches(0.5)
            )
            cf = cap_box.text_frame
            cp = cf.paragraphs[0]
            cp.text = caption
            cp.font.size = SambaSizes.CAPTION
            cp.font.color.rgb = text_color
            cp.font.name = self._get_body_font()
            cp.alignment = PP_ALIGN.CENTER

        return slide

    def add_two_column_slide(self, headline, left_items, right_items,
                              left_header=None, right_header=None, theme=None):
        """
        Add a two-column comparison slide.

        Args:
            headline: Main slide headline
            left_items: List of items for left column
            right_items: List of items for right column
            left_header: Optional header for left column
            right_header: Optional header for right column
            theme: Override default theme
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        # Headline
        head_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, SambaSizes.MARGIN,
            SambaSizes.CONTENT_WIDTH, Inches(1)
        )
        hf = head_box.text_frame
        hp = hf.paragraphs[0]
        hp.text = headline
        hp.font.size = SambaSizes.HEADLINE
        hp.font.color.rgb = text_color
        hp.font.name = self._get_headline_font()

        # Column positions
        left_x = SambaSizes.MARGIN
        right_x = Inches(7)
        col_width = Inches(5.5)
        header_y = Inches(2)
        content_y = Inches(2.75)

        # Left column header
        if left_header:
            lh_box = slide.shapes.add_textbox(left_x, header_y, col_width, Inches(0.5))
            lhf = lh_box.text_frame
            lhp = lhf.paragraphs[0]
            lhp.text = left_header.upper()
            lhp.font.size = SambaSizes.MICRO
            lhp.font.color.rgb = SambaColors.ACID_GREEN
            lhp.font.name = self._get_body_font()
            lhp.font.bold = True

        # Left column content
        lb_box = slide.shapes.add_textbox(left_x, content_y, col_width, Inches(4))
        lbf = lb_box.text_frame
        lbf.word_wrap = True
        for i, item in enumerate(left_items):
            p = lbf.paragraphs[0] if i == 0 else lbf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = SambaSizes.BODY
            p.font.color.rgb = text_color
            p.font.name = self._get_body_font()
            p.space_after = Pt(10)

        # Right column header
        if right_header:
            rh_box = slide.shapes.add_textbox(right_x, header_y, col_width, Inches(0.5))
            rhf = rh_box.text_frame
            rhp = rhf.paragraphs[0]
            rhp.text = right_header.upper()
            rhp.font.size = SambaSizes.MICRO
            rhp.font.color.rgb = SambaColors.ACID_GREEN
            rhp.font.name = self._get_body_font()
            rhp.font.bold = True

        # Right column content
        rb_box = slide.shapes.add_textbox(right_x, content_y, col_width, Inches(4))
        rbf = rb_box.text_frame
        rbf.word_wrap = True
        for i, item in enumerate(right_items):
            p = rbf.paragraphs[0] if i == 0 else rbf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = SambaSizes.BODY
            p.font.color.rgb = text_color
            p.font.name = self._get_body_font()
            p.space_after = Pt(10)

        return slide

    def add_agenda_slide(self, title, items, theme=None):
        """
        Add an agenda/table of contents slide.

        Args:
            title: Slide title (e.g., "Agenda", "Overview")
            items: List of agenda items (strings or dicts with 'text' and optional 'number')
            theme: Override default theme
        """
        slide = self._add_blank_slide(theme)
        _, text_color = self._get_colors(theme)

        # Title
        title_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, SambaSizes.MARGIN,
            SambaSizes.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        tp = tf.paragraphs[0]
        tp.text = title
        tp.font.size = SambaSizes.TITLE_LARGE
        tp.font.color.rgb = text_color
        tp.font.name = self._get_headline_font()

        # Agenda items
        start_y = Inches(2.5)
        item_height = Inches(0.7)

        for i, item in enumerate(items):
            item_text = item if isinstance(item, str) else item.get('text', '')
            item_num = str(i + 1) if isinstance(item, str) else item.get('number', str(i + 1))

            # Number (in Acid Green)
            num_box = slide.shapes.add_textbox(
                SambaSizes.MARGIN, start_y + (i * item_height),
                Inches(0.75), Inches(0.5)
            )
            nf = num_box.text_frame
            np = nf.paragraphs[0]
            np.text = f"{item_num}."
            np.font.size = SambaSizes.SUBHEAD
            np.font.color.rgb = SambaColors.ACID_GREEN
            np.font.name = self._get_headline_font()

            # Item text
            item_box = slide.shapes.add_textbox(
                Inches(1.5), start_y + (i * item_height),
                Inches(10), Inches(0.5)
            )
            itf = item_box.text_frame
            itp = itf.paragraphs[0]
            itp.text = item_text
            itp.font.size = SambaSizes.SUBHEAD
            itp.font.color.rgb = text_color
            itp.font.name = self._get_body_font()

        return slide

    def add_table_slide(self, title, headers, rows, theme=None):
        """
        Add a data table slide with Samba branding.

        Args:
            title: Slide title
            headers: List of column headers
            rows: List of row data (each row is a list of cell values)
            theme: Override default theme
        """
        slide = self._add_blank_slide(theme)
        bg_color, text_color = self._get_colors(theme)

        # Title
        title_box = slide.shapes.add_textbox(
            SambaSizes.MARGIN, SambaSizes.MARGIN,
            SambaSizes.CONTENT_WIDTH, Inches(1)
        )
        tf = title_box.text_frame
        tp = tf.paragraphs[0]
        tp.text = title
        tp.font.size = SambaSizes.HEADLINE
        tp.font.color.rgb = text_color
        tp.font.name = self._get_headline_font()

        # Table dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header
        table_width = Inches(11.5)
        table_height = Inches(min(5, 0.5 * num_rows))
        col_width = table_width / num_cols

        # Create table
        table = slide.shapes.add_table(
            num_rows, num_cols,
            SambaSizes.MARGIN, Inches(2),
            table_width, table_height
        ).table

        # Style header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SambaColors.ACID_GREEN

            # Header text formatting
            para = cell.text_frame.paragraphs[0]
            para.font.size = SambaSizes.TABLE_HEADER
            para.font.color.rgb = SambaColors.VOID_BLACK
            para.font.bold = True
            para.font.name = self._get_body_font()

        # Style data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    SambaColors.DARK_GRAY if (theme or self.theme) == "dark"
                    else SambaColors.LIGHT_GRAY
                )

                para = cell.text_frame.paragraphs[0]
                para.font.size = SambaSizes.TABLE_CELL
                para.font.color.rgb = text_color
                para.font.name = self._get_body_font()

        return slide

    # =========================================================================
    # HUD OVERLAYS ("Tek on Img" - Signature Samba Visual Element)
    # =========================================================================

    def add_hud_coordinates(self, slide, x_ratio, y_ratio, text=None, theme=None):
        """
        Add HUD-style coordinate label to a slide.

        Args:
            slide: The slide to add the overlay to
            x_ratio: X position as ratio of slide width (0.0 to 1.0)
            y_ratio: Y position as ratio of slide height (0.0 to 1.0)
            text: Custom text (default: auto-generates coordinates)
            theme: Override default theme for text color
        """
        _, text_color = self._get_colors(theme)

        # Calculate actual pixel-equivalent coordinates for display
        if text is None:
            display_x = int(x_ratio * 1920)  # Assuming 1080p reference
            display_y = int(y_ratio * 1080)
            text = f"X: {display_x}  Y: {display_y}"

        # Position on slide
        left = Inches(x_ratio * 13.333)
        top = Inches(y_ratio * 7.5)

        coord_box = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.3))
        cf = coord_box.text_frame
        cp = cf.paragraphs[0]
        cp.text = text
        cp.font.size = SambaSizes.HUD_LABEL
        cp.font.color.rgb = SambaColors.ACID_GREEN
        cp.font.name = SambaFonts.get_hud_font()

        return coord_box

    def add_hud_timestamp(self, slide, x_ratio, y_ratio, timestamp=None, theme=None):
        """
        Add HUD-style timestamp to a slide.

        Args:
            slide: The slide to add the overlay to
            x_ratio: X position as ratio of slide width (0.0 to 1.0)
            y_ratio: Y position as ratio of slide height (0.0 to 1.0)
            timestamp: Custom timestamp string (default: current time HH:MM:SS)
            theme: Override default theme
        """
        if timestamp is None:
            timestamp = datetime.now().strftime("%H:%M:%S")

        left = Inches(x_ratio * 13.333)
        top = Inches(y_ratio * 7.5)

        ts_box = slide.shapes.add_textbox(left, top, Inches(1.5), Inches(0.3))
        tf = ts_box.text_frame
        tp = tf.paragraphs[0]
        tp.text = timestamp
        tp.font.size = SambaSizes.HUD_LABEL
        tp.font.color.rgb = SambaColors.ACID_GREEN
        tp.font.name = SambaFonts.get_hud_font()

        return ts_box

    def add_tracking_dot(self, slide, x_ratio, y_ratio, label=None, color=None):
        """
        Add a tracking dot overlay (for subject tracking visualization).

        Args:
            slide: The slide to add the overlay to
            x_ratio: X position as ratio of slide width (0.0 to 1.0)
            y_ratio: Y position as ratio of slide height (0.0 to 1.0)
            label: Optional label text above the dot
            color: Dot color (default: Acid Green)
        """
        dot_color = color or SambaColors.ACID_GREEN

        # Calculate position
        left = Inches(x_ratio * 13.333) - (SambaSizes.HUD_DOT_SIZE / 2)
        top = Inches(y_ratio * 7.5) - (SambaSizes.HUD_DOT_SIZE / 2)

        # Add dot (oval shape)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top,
            SambaSizes.HUD_DOT_SIZE, SambaSizes.HUD_DOT_SIZE
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()

        # Add label above dot
        if label:
            label_box = slide.shapes.add_textbox(
                left - Inches(0.5), top - Inches(0.25),
                Inches(1.2), Inches(0.2)
            )
            lf = label_box.text_frame
            lp = lf.paragraphs[0]
            lp.text = label.upper()
            lp.font.size = SambaSizes.HUD_LABEL
            lp.font.color.rgb = dot_color
            lp.font.name = SambaFonts.get_hud_font()
            lp.alignment = PP_ALIGN.CENTER

        return dot

    def add_bounding_box(self, slide, x_ratio, y_ratio, width_ratio, height_ratio,
                         label=None, color=None):
        """
        Add a bounding box overlay (for object detection visualization).

        Args:
            slide: The slide to add the overlay to
            x_ratio: X position of top-left corner (0.0 to 1.0)
            y_ratio: Y position of top-left corner (0.0 to 1.0)
            width_ratio: Width as ratio of slide width
            height_ratio: Height as ratio of slide height
            label: Optional label text (displayed at top of box)
            color: Box color (default: Acid Green)
        """
        box_color = color or SambaColors.ACID_GREEN

        # Calculate position and size
        left = Inches(x_ratio * 13.333)
        top = Inches(y_ratio * 7.5)
        width = Inches(width_ratio * 13.333)
        height = Inches(height_ratio * 7.5)

        # Add rectangle (no fill, just outline)
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        box.fill.background()  # Transparent fill
        box.line.color.rgb = box_color
        box.line.width = Pt(2)

        # Add label at top of box
        if label:
            label_box = slide.shapes.add_textbox(
                left, top - Inches(0.25),
                width, Inches(0.2)
            )
            lf = label_box.text_frame
            lp = lf.paragraphs[0]
            lp.text = label.upper()
            lp.font.size = SambaSizes.HUD_LABEL
            lp.font.color.rgb = box_color
            lp.font.name = SambaFonts.get_hud_font()

        return box

    def add_wireframe_grid(self, slide, rows=4, cols=6, color=None, opacity=0.3):
        """
        Add a subtle wireframe grid overlay to a slide.

        Args:
            slide: The slide to add the overlay to
            rows: Number of horizontal grid lines
            cols: Number of vertical grid lines
            color: Grid line color (default: white with opacity)
            opacity: Line opacity (0.0 to 1.0) - Note: limited support in PPTX
        """
        grid_color = color or SambaColors.STARK_WHITE

        slide_w = 13.333
        slide_h = 7.5

        # Horizontal lines
        for i in range(1, rows):
            y = (i / rows) * slide_h
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(y),
                Inches(slide_w), Pt(0.5)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = grid_color
            line.line.fill.background()

        # Vertical lines
        for i in range(1, cols):
            x = (i / cols) * slide_w
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(0),
                Pt(0.5), Inches(slide_h)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = grid_color
            line.line.fill.background()

        return slide

    def add_hud_label(self, slide, x_ratio, y_ratio, text, color=None):
        """
        Add a generic HUD-style label anywhere on a slide.

        Args:
            slide: The slide to add the label to
            x_ratio: X position as ratio (0.0 to 1.0)
            y_ratio: Y position as ratio (0.0 to 1.0)
            text: Label text
            color: Text color (default: Acid Green)
        """
        label_color = color or SambaColors.ACID_GREEN

        left = Inches(x_ratio * 13.333)
        top = Inches(y_ratio * 7.5)

        label_box = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.25))
        lf = label_box.text_frame
        lp = lf.paragraphs[0]
        lp.text = text.upper()
        lp.font.size = SambaSizes.HUD_LABEL
        lp.font.color.rgb = label_color
        lp.font.name = SambaFonts.get_hud_font()

        return label_box

    # =========================================================================
    # UTILITIES
    # =========================================================================

    def add_slide_notes(self, slide, notes_text):
        """
        Add speaker notes to a slide.

        Args:
            slide: The slide to add notes to
            notes_text: The speaker notes text
        """
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = notes_text
        return notes_slide

    def add_logo_to_slide(self, slide, position="bottom-right", size="small", theme=None):
        """
        Add Samba logo to an existing slide.

        Args:
            slide: The slide object
            position: "bottom-right", "bottom-left", "top-right", "top-left", "center"
            size: "small" (1.5in), "medium" (2.5in), "large" (4in)
            theme: Override theme for logo color selection (uses slide/instance theme if None)
        """
        sizes = {"small": Inches(1.5), "medium": Inches(2.5), "large": Inches(4)}
        width = sizes.get(size, sizes["small"])

        # Use provided theme or fall back to instance theme
        effective_theme = theme or self.theme

        # Get logo path from LogoManager (context-aware)
        logo_path = self.logo_manager.get_logo_path(self.context, effective_theme)

        if not logo_path.exists():
            # Fallback to direct file check
            logo_file = "samba-logo-white.png" if effective_theme == "dark" else "samba-logo-black.png"
            logo_path = self.logos_dir / logo_file
            if not logo_path.exists():
                return None

        # Calculate position
        positions = {
            "bottom-right": (Inches(13.333) - width - SambaSizes.MARGIN, Inches(6.5)),
            "bottom-left": (SambaSizes.MARGIN, Inches(6.5)),
            "top-right": (Inches(13.333) - width - SambaSizes.MARGIN, SambaSizes.MARGIN),
            "top-left": (SambaSizes.MARGIN, SambaSizes.MARGIN),
            "center": ((Inches(13.333) - width) / 2, Inches(5.5)),  # Centered lower
        }
        left, top = positions.get(position, positions["bottom-right"])

        return slide.shapes.add_picture(str(logo_path), left, top, width=width)

    def save(self, output_path):
        """Save the presentation to a file."""
        self.prs.save(output_path)
        return output_path


# =============================================================================
# CLI USAGE
# =============================================================================

if __name__ == "__main__":
    # Demo: Create a comprehensive sample presentation
    prs = SambaPresentation(theme="dark")

    # Title slide
    prs.add_title_slide(
        "AI-Driven Media Intelligence",
        "Samba TV"
    )

    # Agenda slide (NEW)
    prs.add_agenda_slide("Agenda", [
        "The Challenge",
        "Our Solution",
        "Key Metrics",
        "Platform Comparison",
        "Next Steps"
    ])

    # Section divider (now centered by default)
    prs.add_section_divider("The Problem")

    prs.add_content_slide(
        "Traditional measurement is broken",
        [
            "Fragmented data across platforms",
            "No unified cross-device view",
            "Delayed insights miss the moment",
            "Walled gardens limit visibility"
        ]
    )

    prs.add_section_divider("The Solution")

    # Metrics slide
    prs.add_metrics_slide([
        {"value": "500M+", "label": "Connected Devices", "accent": True},
        {"value": "24/7", "label": "Real-Time Data"},
        {"value": "99.9%", "label": "Measurement Accuracy"}
    ])

    # Two-column comparison slide (NEW)
    prs.add_two_column_slide(
        "Platform Comparison",
        left_items=[
            "Panel-based sampling",
            "Delayed reporting (24-48h)",
            "Limited cross-device tracking",
            "Siloed data sources"
        ],
        right_items=[
            "Census-level measurement",
            "Real-time insights",
            "True cross-device identity",
            "Unified data platform"
        ],
        left_header="Traditional",
        right_header="Samba TV"
    )

    # Table slide (NEW)
    prs.add_table_slide(
        "Quarterly Performance",
        headers=["Metric", "Q3 2024", "Q4 2024", "Change"],
        rows=[
            ["Active Devices", "485M", "512M", "+5.6%"],
            ["Data Points/Day", "2.1B", "2.8B", "+33%"],
            ["Avg. Response Time", "45ms", "38ms", "-15%"],
        ]
    )

    prs.add_content_slide(
        "Unified media intelligence",
        [
            "Cross-platform attribution",
            "Real-time audience insights",
            "AI-powered optimization",
            "Privacy-first approach"
        ]
    )

    # Quote slide
    prs.add_quote_slide(
        "Samba TV transformed how we understand our audience across every screen.",
        "Fortune 500 CMO"
    )

    # Demo: HUD overlays on an image slide (NEW)
    # Note: This would normally use a real image path
    img_slide = prs.add_image_slide(
        None,  # Would be: "path/to/lifestyle_photo.jpg"
        caption="AI-powered audience identification"
    )

    # Add "Tek on Img" HUD overlays
    prs.add_hud_coordinates(img_slide, 0.85, 0.05)
    prs.add_hud_timestamp(img_slide, 0.85, 0.1)
    prs.add_tracking_dot(img_slide, 0.35, 0.4, label="VIEWER A")
    prs.add_tracking_dot(img_slide, 0.55, 0.45, label="VIEWER B")
    prs.add_bounding_box(img_slide, 0.3, 0.25, 0.4, 0.5, label="LIVING ROOM")
    prs.add_hud_label(img_slide, 0.02, 0.95, "SAMBA VISION v2.1")

    # Add speaker notes to the image slide
    prs.add_slide_notes(img_slide,
        "This slide demonstrates our AI-powered viewer identification. "
        "The HUD overlays show real-time tracking coordinates and subject identification. "
        "Key talking point: Privacy-preserving identification at scale."
    )

    # End slide
    prs.add_end_slide()

    # =========================================================================
    # NEW: PRIMITIVES DEMO - Custom layouts with flexible positioning
    # =========================================================================

    # Demo slide using primitives for custom layout
    custom_slide = prs.add_slide()  # New public method

    # Add accent bar on left edge
    prs.add_accent_bar(custom_slide, position="left", thickness=0.01, length=0.6)

    # Positioned headline
    prs.add_headline(custom_slide, "Custom Layout Demo", x=0.08, y=0.1, width=0.5,
                     size="medium", color="text")

    # Positioned body text
    prs.add_body_text(custom_slide, "This slide was built using the new primitives API. "
                      "Each element can be positioned anywhere on the slide.",
                      x=0.08, y=0.25, width=0.45, height=0.2, size="body")

    # Bullet list in left column
    prs.add_bullets(custom_slide, [
        "Flexible x/y positioning",
        "Multiple coordinate systems",
        "Brand-compliant styling",
        "Easy element grouping"
    ], x=0.08, y=0.45, width=0.4)

    # Metrics in right column using primitives
    prs.add_metric(custom_slide, "100%", "FLEXIBILITY", x=0.55, y=0.15, accent=True)
    prs.add_metric(custom_slide, "3", "COORD SYSTEMS", x=0.75, y=0.15)

    # Shape primitives
    prs.add_rectangle(custom_slide, x=0.55, y=0.55, width=0.4, height=0.3,
                      fill="black", border=True, border_color="accent")
    prs.add_circle(custom_slide, x=0.65, y=0.7, size=0.08, fill="accent")
    prs.add_circle(custom_slide, x=0.85, y=0.7, size=0.08, fill="white")

    # Horizontal line
    prs.add_line(custom_slide, start=(0.55, 0.52), end=(0.95, 0.52),
                 color="gray", thickness=1)

    # Demo: Element grouping
    group = prs.create_group(custom_slide, "footer_group")
    footer_rect = prs.add_rectangle(custom_slide, x=0.05, y=0.92, width=0.9, height=0.05,
                                    fill="black")
    footer_text = prs.add_body_text(custom_slide, "Built with Samba Slides Primitives API",
                                    x=0.05, y=0.92, width=0.9, height=0.05,
                                    size="micro", color="accent", align="center")
    prs.add_to_group(group, footer_rect)
    prs.add_to_group(group, footer_text)

    output = prs.save("samba_demo.pptx")
    print(f"Created: {output}")
    print("\nFeatures demonstrated:")
    print("- Agenda slide with numbered items")
    print("- Two-column comparison slide")
    print("- Data table slide")
    print("- Full-bleed image slide with caption")
    print("- HUD overlays (coordinates, timestamp, tracking dots, bounding box)")
    print("- Speaker notes")
    print("\nNEW Primitives API:")
    print("- add_slide() - Create blank themed slide")
    print("- add_headline() - Positioned headline text")
    print("- add_body_text() - Positioned body text")
    print("- add_bullets() - Bullet list")
    print("- add_metric() - Value/label metric pair")
    print("- add_rectangle() - Colored rectangle")
    print("- add_accent_bar() - Brand accent bar")
    print("- add_circle() - Circle shape")
    print("- add_line() - Line shape")
    print("- create_group() / add_to_group() - Logical element grouping")
